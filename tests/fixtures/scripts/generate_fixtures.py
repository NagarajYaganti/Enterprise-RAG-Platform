"""Generates real, small sample files for each supported ingestion format.

Run with: uv run --group fixtures python tests/fixtures/scripts/generate_fixtures.py

Every fixture is real content produced by these libraries at generation
time — not a hand-fabricated stand-in for a real file.
"""

from email.message import EmailMessage
from pathlib import Path

from docx import Document as DocxDocument
from fpdf import FPDF
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pypdf import PdfReader, PdfWriter

OUTPUT_DIR = Path(__file__).resolve().parents[1] / "documents"

PDF_TEXT = "Quarterly Compliance Report"
DOCX_HEADING = "Lending Policy Overview"
DOCX_BODY = "All loan applications must be reviewed within five business days."
PPTX_TITLE = "Q3 Retail Sales Summary"
PPTX_BODY = "Revenue grew twelve percent quarter over quarter."
XLSX_HEADERS = ["product", "units_sold", "revenue"]
XLSX_ROW = ["widget-a", 42, 1050]
HTML_HEADING = "Onboarding Runbook"
HTML_BODY = "Follow these steps to onboard a new enterprise tenant."
OCR_TEXT = "INVOICE NUMBER 48213"
SCANNED_PDF_TEXT = "PURCHASE ORDER 77042"
STT_TEXT = "the quarterly earnings call starts at nine a m eastern time"
EML_SUBJECT = "Contract renewal reminder"
EML_BODY = "Please review the attached renewal terms before Friday."
TXT_TEXT = "Standard support response times are documented in the service level agreement."
MD_TEXT = "# Runbook\n\nRestart the ingestion worker if the queue depth exceeds one thousand."
CSV_TEXT = "product,units_sold,revenue\nwidget-b,17,425\n"
JSON_TEXT = '{"policy": "refund", "window_days": 30}'
XML_TEXT = "<policy><name>refund</name><window_days>30</window_days></policy>"
# A real legacy (non-UTF-8) encoding, to exercise PlainTextParser's
# charset-detection fallback path for real, not just the common UTF-8 case.
LEGACY_TEXT = "Le délai de remboursement standard est de trente jours ouvrés."

# Real sentences (not transliterated placeholders) in each of lingua's
# supported non-Latin languages (preprocessing/language_detect.py), all
# expressing the same underlying policy statement as DOCX_BODY above so the
# multilingual integration test can sanity-check meaning stayed on-topic,
# not just that a script was detected.
ARABIC_TEXT = "يجب مراجعة جميع طلبات القروض خلال خمسة أيام عمل."
CHINESE_TEXT = "所有贷款申请必须在五个工作日内完成审核。"
HINDI_TEXT = "सभी ऋण आवेदनों की समीक्षा पांच कार्य दिवसों के भीतर की जानी चाहिए।"
MIXED_EN_HEADING = "Loan Policy"
MIXED_EN_BODY = "All loan applications must be reviewed within five business days."
MIXED_AR_HEADING = "سياسة القروض"


def generate_pdf() -> None:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=16)
    pdf.cell(text=PDF_TEXT)
    pdf.output(str(OUTPUT_DIR / "sample.pdf"))


def generate_docx() -> None:
    doc = DocxDocument()
    doc.add_heading(DOCX_HEADING, level=1)
    doc.add_paragraph(DOCX_BODY)
    doc.save(str(OUTPUT_DIR / "sample.docx"))


def generate_pptx() -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = PPTX_TITLE
    body = slide.placeholders[1]
    body.text_frame.text = PPTX_BODY
    prs.save(str(OUTPUT_DIR / "sample.pptx"))


def generate_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.append(XLSX_HEADERS)
    ws.append(XLSX_ROW)
    wb.save(str(OUTPUT_DIR / "sample.xlsx"))


def generate_html() -> None:
    html = f"<html><body><h1>{HTML_HEADING}</h1><p>{HTML_BODY}</p></body></html>"
    (OUTPUT_DIR / "sample.html").write_text(html)


def generate_ocr_image() -> None:
    img = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype(
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", size=28
    )
    draw.text((10, 30), OCR_TEXT, fill="black", font=font)
    img.save(str(OUTPUT_DIR / "sample_ocr.png"))


def generate_scanned_pdf() -> None:
    # A genuinely image-only PDF -- no embedded text layer at all. Pillow
    # renders the text onto a raster image, then saves that image directly
    # as a single-page PDF, exactly the shape a real flatbed-scanned
    # document has. Distinct from sample_ocr.png (a raw PNG, exercising
    # ParserPolicy's separate images_via_ocr mime-type route): this exists
    # to exercise the unstructured-library PDF route's OWN internal
    # image-only-PDF/OCR fallback, invisible to ParserPolicy itself (which
    # routes every application/pdf, native or scanned, to the same
    # "unstructured" route -- see libs/connectors/parser_policy.py).
    img = Image.new("RGB", (600, 150), color="white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype(
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", size=28
    )
    draw.text((10, 50), SCANNED_PDF_TEXT, fill="black", font=font)
    img.save(str(OUTPUT_DIR / "sample_scanned.pdf"), "PDF")


def generate_audio() -> None:
    import pyttsx3

    engine = pyttsx3.init()
    # Default espeak-ng rate (200wpm) transcribes poorly with faster-whisper on
    # this synthetic voice; slowing it down measurably improves accuracy
    # (verified: "quarterly earnings call" transcribes correctly at rate=130
    # but not at the default rate).
    engine.setProperty("rate", 130)
    engine.save_to_file(STT_TEXT, str(OUTPUT_DIR / "sample_audio.wav"))
    engine.runAndWait()


def generate_eml() -> None:
    msg = EmailMessage()
    msg["Subject"] = EML_SUBJECT
    msg["From"] = "sender@example.com"
    msg["To"] = "recipient@example.com"
    msg.set_content(EML_BODY)
    (OUTPUT_DIR / "sample.eml").write_bytes(msg.as_bytes())


def generate_encrypted_pdf() -> None:
    # Real encryption via pypdf.PdfWriter.encrypt on the real sample.pdf
    # generated above -- not a fabricated stand-in.
    reader = PdfReader(str(OUTPUT_DIR / "sample.pdf"))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password="test-password-123")
    with (OUTPUT_DIR / "sample_encrypted.pdf").open("wb") as f:
        writer.write(f)


def generate_plain_text_fixtures() -> None:
    (OUTPUT_DIR / "sample.txt").write_text(TXT_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample.md").write_text(MD_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample.csv").write_text(CSV_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample.json").write_text(JSON_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample.xml").write_text(XML_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample_legacy_encoding.txt").write_text(LEGACY_TEXT, encoding="windows-1252")


def generate_multilingual_fixtures() -> None:
    (OUTPUT_DIR / "sample_arabic.txt").write_text(ARABIC_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample_chinese.txt").write_text(CHINESE_TEXT, encoding="utf-8")
    (OUTPUT_DIR / "sample_hindi.txt").write_text(HINDI_TEXT, encoding="utf-8")
    # An English section followed by an Arabic section, in one document --
    # proves per-section (not per-document) language detection, since the
    # document-wide majority-language vote would otherwise mislabel one of
    # the two sections.
    html = (
        "<html><body>"
        f"<h1>{MIXED_EN_HEADING}</h1><p>{MIXED_EN_BODY}</p>"
        f"<h1>{MIXED_AR_HEADING}</h1><p>{ARABIC_TEXT}</p>"
        "</body></html>"
    )
    (OUTPUT_DIR / "sample_mixed_en_ar.html").write_text(html, encoding="utf-8")


def generate_corrupt_pdf() -> None:
    # Genuinely not a valid PDF at all (no %PDF header, no xref table) --
    # real corruption, not a fabricated-but-plausible file.
    (OUTPUT_DIR / "sample_corrupt.pdf").write_bytes(
        b"this is not a real pdf file, just garbage bytes for a corruption test"
    )


if __name__ == "__main__":
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    generate_pdf()
    generate_docx()
    generate_pptx()
    generate_xlsx()
    generate_html()
    generate_ocr_image()
    generate_scanned_pdf()
    generate_audio()
    generate_eml()
    generate_encrypted_pdf()
    generate_corrupt_pdf()
    generate_plain_text_fixtures()
    generate_multilingual_fixtures()
    print(f"Generated fixtures in {OUTPUT_DIR}")
